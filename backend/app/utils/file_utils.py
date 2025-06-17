import os
import fitz  # PyMuPDF
import docx
import pptx
import textract

def save_file(file, destination: str):
    """ Save file to a specified destination """
    with open(destination, "wb") as f:
        f.write(file)
        
def remove_temp_file(file_path: str):
    """ Remove temporary file after processing """
    try:
        os.remove(file_path)
    except Exception as e:
        print(f"Error removing temp file: {str(e)}")

def extract_text_from_pdf(file_path):
    """Extract text from a PDF file."""
    doc = fitz.open(file_path)
    return "\n".join([page.get_text() for page in doc])

def extract_text_from_word(file_path):
    """Extract text from a Word document."""
    doc = docx.Document(file_path)
    return "\n".join([para.text for para in doc.paragraphs])

def extract_text_from_ppt(file_path):
    """Extract text from a PowerPoint presentation."""
    prs = pptx.Presentation(file_path)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text)

def extract_text_generic(file_path):
    """Extract text from other file formats using textract."""
    try:
        return textract.process(file_path).decode("utf-8")
    except Exception as e:
        return f"Error extracting text: {str(e)}"

def extract_text(file_path: str, file_ext: str) -> str:
    """
    Extracts text from a file based on its extension.
    """
    if file_ext == "pdf":
        return extract_text_from_pdf(file_path)
    elif file_ext in ["doc", "docx"]:
        return extract_text_from_word(file_path)
    elif file_ext in ["ppt", "pptx"]:
        return extract_text_from_ppt(file_path)
    else:
        return extract_text_generic(file_path) 