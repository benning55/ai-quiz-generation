from fastapi import FastAPI, UploadFile, File, Depends, Request, HTTPException, Query
from sqlalchemy.orm import Session
from sqlalchemy import func
from typing import List, Optional
from .db import SessionLocal, engine
from . import models
from .groq_client import generate_questions
import os
import random
from pydantic import BaseModel, Field
import jwt
from fastapi.security import HTTPBearer, HTTPAuthorizationCredentials
import hmac
import hashlib
import json
from fastapi.middleware.cors import CORSMiddleware
from dotenv import load_dotenv
import stripe
from datetime import datetime, timedelta
from .models import Base, User, Flashcard, Payment

# Import the FastAPI app from __init__.py
from . import app

# Load environment variables
load_dotenv()

# Initialize Stripe
stripe.api_key = os.getenv("STRIPE_SECRET_KEY")

# Dependency to get the database session
def get_db():
    db = SessionLocal()  # Create a new database session
    try:
        yield db  # This will yield a database session to be used in the route
    finally:
        db.close()  # Ensure that the session is closed after use

# Extract text from PDF
def extract_text_from_pdf(file_path):
    import fitz  # PyMuPDF for PDFs
    doc = fitz.open(file_path)
    return "\n".join([page.get_text() for page in doc])

# Extract text from Word
def extract_text_from_word(file_path):
    import docx  # python-docx for Word
    doc = docx.Document(file_path)
    return "\n".join([para.text for para in doc.paragraphs])

# Extract text from PowerPoint
def extract_text_from_ppt(file_path):
    import pptx  # python-pptx for PowerPoint
    prs = pptx.Presentation(file_path)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text)

# General extractor for other formats
def extract_text_generic(file_path):
    import textract  # For various formats
    try:
        return textract.process(file_path).decode("utf-8")
    except Exception as e:
        return f"Error extracting text: {str(e)}"

@app.post("/extract-text")
@app.post("/extract-text/")
async def extract_text(file: UploadFile = File(...), db: Session = Depends(get_db)):
    file_ext = file.filename.split(".")[-1].lower()
    file_path = f"temp.{file_ext}"

    # Save uploaded file temporarily
    with open(file_path, "wb") as f:
        f.write(await file.read())

    # Extract text based on file type
    if file_ext == "pdf":
        text = extract_text_from_pdf(file_path)
    elif file_ext in ["doc", "docx"]:
        text = extract_text_from_word(file_path)
    elif file_ext in ["ppt", "pptx"]:
        text = extract_text_from_ppt(file_path)
    else:
        text = extract_text_generic(file_path)

    # Optionally, save the document to the database
    db.add(models.Document(name=file.filename, content=text))
    db.commit()

    # Clean up the temporary file
    os.remove(file_path)

    quiz_json = generate_questions(text)
    
    # Return the extracted text
    return {"filename": file.filename, "extracted_text": text, 'quiz': quiz_json}

@app.get("/")
async def hello(request: Request):
    # Access information from the request object
    client_ip = request.client.host  # Get client IP address
    headers = request.headers  # Access headers
    
    # You can also access query parameters, cookies, etc.
    
    return {
        "message": "Hello",
        "client_ip": client_ip,
        "headers": dict(headers)
    }

# Add health check endpoint
@app.get("/health")
async def health_check():
    return {"status": "healthy"}

# Pydantic models for API request/response
class FlashcardBase(BaseModel):
    question: str
    answer: str
    tags: Optional[List[str]] = None
    category: Optional[str] = None

class FlashcardCreate(FlashcardBase):
    pass

class Flashcard(FlashcardBase):
    id: int
    
    class Config:
        orm_mode = True

class QuizRequest(BaseModel):
    category: Optional[str] = None
    count: int = 10
    question_types: List[str] = ["multiple_choice", "true_false"]

# Routes for flashcards
@app.get("/flashcards/", response_model=List[Flashcard])
def get_flashcards(
    skip: int = 0, 
    limit: int = 100, 
    category: Optional[str] = None,
    tag: Optional[str] = None,
    db: Session = Depends(get_db)
):
    """Get all flashcards with optional filtering by category or tag"""
    query = db.query(models.Flashcard)
    
    if category:
        query = query.filter(models.Flashcard.category == category)
        
    if tag:
        query = query.filter(models.Flashcard.tags.contains([tag]))
        
    flashcards = query.offset(skip).limit(limit).all()
    return flashcards

@app.get("/flashcards/{flashcard_id}", response_model=Flashcard)
def get_flashcard(flashcard_id: int, db: Session = Depends(get_db)):
    """Get a specific flashcard by ID"""
    flashcard = db.query(models.Flashcard).filter(models.Flashcard.id == flashcard_id).first()
    if flashcard is None:
        raise HTTPException(status_code=404, detail="Flashcard not found")
    return flashcard

@app.post("/flashcards/", response_model=Flashcard)
def create_flashcard(flashcard: FlashcardCreate, db: Session = Depends(get_db)):
    """Create a new flashcard"""
    db_flashcard = models.Flashcard(**flashcard.dict())
    db.add(db_flashcard)
    db.commit()
    db.refresh(db_flashcard)
    return db_flashcard

# Initialize the security scheme
security = HTTPBearer(auto_error=False)

# Authentication utilities
async def get_user_id(credentials: Optional[HTTPAuthorizationCredentials] = Depends(security)):
    """Verify JWT token and return user ID"""
    if not credentials:
        return None
    
    try:
        token = credentials.credentials
        # For Clerk, you'd normally validate with Clerk's API
        # This is a simplified version that just checks if the token exists
        # In production, you should verify with Clerk's JWT verification
        return True  # Return user is authenticated
    except Exception as e:
        return None

# User management models
class UserBase(BaseModel):
    email: Optional[str] = None
    first_name: Optional[str] = None
    last_name: Optional[str] = None
    image_url: Optional[str] = None

class UserCreate(UserBase):
    clerk_id: str

class UserResponse(UserBase):
    id: int
    clerk_id: str
    created_at: Optional[str] = None
    
    class Config:
        orm_mode = True
        
    @classmethod
    def from_orm(cls, obj):
        # Create a dict from the ORM model instance
        dict_obj = {col.name: getattr(obj, col.name) for col in obj.__table__.columns}
        
        # Convert datetime to string
        if dict_obj.get('created_at'):
            dict_obj['created_at'] = dict_obj['created_at'].isoformat()
        if dict_obj.get('updated_at'):
            dict_obj['updated_at'] = dict_obj['updated_at'].isoformat()
            
        # Create an instance of this model using the dict
        return cls(**dict_obj)

# Function to get user from token
async def get_current_user(
    credentials: Optional[HTTPAuthorizationCredentials] = Depends(security),
    db: Session = Depends(get_db),
    request: Request = None
):
    print("==== Authentication Debug ====")
    
    # Check for token in Authorization header
    token = None
    if credentials:
        token = credentials.credentials
        print("Token found in Authorization header")
    
    # If no token in header, check for Clerk session cookie
    if not token and request:
        cookies = request.cookies
        session_cookie = cookies.get("__session")
        if session_cookie:
            token = session_cookie
            print("Token found in Clerk session cookie")
    
    if not token:
        print("No authentication token found")
        return None
    
    try:
        # Decode and verify the JWT token
        try:
            decoded_token = jwt.decode(
                token,
                options={"verify_signature": False}  # In production, should verify with Clerk's public key
            )
            
            # Print token details for debugging
            print(f"Decoded token: {decoded_token}")
            
            # Extract the clerk_id from the token
            clerk_id = decoded_token.get("sub")
            print(f"Extracted clerk_id: {clerk_id}")
            
            if not clerk_id:
                print("No clerk_id found in token")
                return None
            
            # Find user in database
            user = db.query(models.User).filter(models.User.clerk_id == clerk_id).first()
            print(f"Found user in DB: {user is not None}")
            
            # If user doesn't exist yet, create one (auto-registration)
            if not user:
                print("Creating new user in database")
                user = models.User(
                    clerk_id=clerk_id,
                    email=decoded_token.get("email"),
                    first_name=decoded_token.get("given_name", ""),
                    last_name=decoded_token.get("family_name", "")
                )
                db.add(user)
                db.commit()
                db.refresh(user)
                print(f"New user created with ID: {user.id}")
            else:
                print(f"Existing user found with ID: {user.id}")
            
            return user
        except jwt.DecodeError as e:
            print(f"JWT decode error: {str(e)}")
            print(f"Token format issue, trying to handle as plain token: {token[:20]}...")
            # Maybe it's not a JWT format, try using the token as clerk_id directly
            user = db.query(models.User).filter(models.User.clerk_id == token).first()
            if user:
                print(f"Found user with token as clerk_id: {user.id}")
                return user
            return None
    except Exception as e:
        print(f"Authentication error: {str(e)}")
        import traceback
        print(traceback.format_exc())
        return None

# User management endpoints
@app.post("/users/", response_model=UserResponse)
async def create_user(user: UserCreate, db: Session = Depends(get_db)):
    """Create a new user (called when a user signs up with Clerk)"""
    # Check if user already exists
    db_user = db.query(models.User).filter(models.User.clerk_id == user.clerk_id).first()
    if db_user:
        return db_user
    
    # Create new user
    db_user = models.User(
        clerk_id=user.clerk_id,
        email=user.email,
        first_name=user.first_name,
        last_name=user.last_name,
        image_url=user.image_url
    )
    db.add(db_user)
    db.commit()
    db.refresh(db_user)
    return db_user

@app.get("/users/me", response_model=UserResponse)
async def get_current_user_info(current_user = Depends(get_current_user)):
    """Get current user info"""
    if not current_user:
        raise HTTPException(status_code=401, detail="Not authenticated")
    return current_user

@app.post("/create-payment-intent")
async def create_payment_intent(
    current_user: User = Depends(get_current_user),
    db: Session = Depends(get_db)
):
    if not current_user:
        raise HTTPException(status_code=401, detail="Not authenticated")
    
    try:
        # Create a PaymentIntent with the order amount and currency
        payment_intent = stripe.PaymentIntent.create(
            amount=2500,  # $25.00 CAD
            currency='cad',
            automatic_payment_methods={
                'enabled': True,
            },
            metadata={
                'user_id': str(current_user.id)
            }
        )
        
        return {"clientSecret": payment_intent.client_secret}
    except Exception as e:
        raise HTTPException(status_code=400, detail=str(e))

@app.post("/webhook")
async def stripe_webhook(request: Request, db: Session = Depends(get_db)):
    payload = await request.body()
    sig_header = request.headers.get('stripe-signature')
    
    try:
        event = stripe.Webhook.construct_event(
            payload, sig_header, os.getenv("STRIPE_WEBHOOK_SECRET")
        )
    except ValueError as e:
        raise HTTPException(status_code=400, detail="Invalid payload")
    except stripe.error.SignatureVerificationError as e:
        raise HTTPException(status_code=400, detail="Invalid signature")

    if event.type == 'payment_intent.succeeded':
        payment_intent = event.data.object
        user_id = payment_intent.metadata.get('user_id')
        
        # Create payment record
        payment = Payment(
            user_id=user_id,
            stripe_payment_intent_id=payment_intent.id,
            amount=payment_intent.amount,
            status='succeeded',
            expires_at=datetime.utcnow() + timedelta(days=30)  # 30-day access
        )
        db.add(payment)
        db.commit()

    return {"status": "success"}

# Modify the quiz generation endpoint
@app.post("/generate-quiz-from-flashcards/")
async def generate_quiz_from_flashcards(
    request: QuizRequest,
    current_user: User = Depends(get_current_user),
    db: Session = Depends(get_db)
):
    try:
        # Check if user is authenticated
        is_authenticated = current_user is not None
        
        # Get user's active payment if authenticated
        active_payment = None
        if is_authenticated:
            active_payment = db.query(Payment).filter(
                Payment.user_id == current_user.id,
                Payment.status == 'succeeded',
                Payment.expires_at > datetime.utcnow()
            ).first()
        
        # Determine question count based on user status
        if not is_authenticated:
            question_count = 3  # Free tier
        elif active_payment:
            question_count = request.count  # Full access
        else:
            question_count = min(request.count, 3)  # Logged in but not paid
            
        # Build query based on request parameters
        query = db.query(models.Flashcard)
        
        if request.category:
            query = query.filter(models.Flashcard.category == request.category)
        
        # Get total count of matching flashcards
        total_cards = query.count()
        print(f"Total available flashcards: {total_cards}")
        
        if total_cards == 0:
            raise HTTPException(status_code=404, detail="No flashcards found matching the criteria")
        
        # Limit the number of flashcards to what was requested
        count = min(question_count, total_cards)
        print(f"Requested flashcard count: {count}")
        
        print(f"Final flashcard count: {count}")
        
        # Get random flashcards
        random_offset = random.randint(0, max(0, total_cards - count))
        selected_flashcards = query.offset(random_offset).limit(count).all()
        
        # Format the flashcards into a study material
        study_material = "\n\n".join([
            f"Question: {card.question}\nAnswer: {card.answer}" 
            for card in selected_flashcards
        ])
        
        # Use groq to generate quiz questions based on flashcards
        quiz_json = generate_questions(
            study_material, 
            question_count=count,
            question_types=request.question_types
        )
        
        return {
            "questions": quiz_json,
            "user_status": {
                "is_authenticated": is_authenticated,
                "has_active_payment": active_payment is not None,
                "question_count": question_count
            }
        }
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

class FlashcardsImport(BaseModel):
    flashcards: List[FlashcardCreate]

@app.post("/import-flashcards/")
async def import_flashcards(data: FlashcardsImport, db: Session = Depends(get_db)):
    """Import multiple flashcards from JSON"""
    imported_flashcards = []
    
    for flashcard_data in data.flashcards:
        db_flashcard = models.Flashcard(**flashcard_data.dict())
        db.add(db_flashcard)
        imported_flashcards.append(db_flashcard)
    
    db.commit()
    
    # Refresh all the flashcards to get their IDs
    for flashcard in imported_flashcards:
        db.refresh(flashcard)
    
    return {
        "success": True,
        "message": f"{len(imported_flashcards)} flashcards imported successfully",
        "imported_flashcards": imported_flashcards
    }

# Clerk webhook models
class ClerkWebhookPayload(BaseModel):
    data: dict = Field(...)
    object: str
    type: str

@app.post("/api/webhooks/clerk")
async def clerk_webhook(request: Request, db: Session = Depends(get_db)):
    """Handle Clerk webhook events (user creation, deletion, etc.)"""
    
    # Get the webhook signature from headers
    clerk_signature = request.headers.get("svix-signature")
    if not clerk_signature:
        raise HTTPException(status_code=400, detail="Missing Clerk signature")
    
    # Get the webhook secret from environment variable
    webhook_secret = os.environ.get("CLERK_WEBHOOK_SECRET")
    if not webhook_secret:
        # For development, you can hardcode a secret (remove in production)
        webhook_secret = "your_webhook_secret_here"
        
    # Get the raw request body
    body = await request.body()
    body_str = body.decode("utf-8")
    
    # Verify the webhook signature
    try:
        # In a real implementation, you would verify the signature with your webhook secret
        # For now, we'll skip this step for demonstration
        # This should be implemented in production
        payload = json.loads(body_str)
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"Invalid payload: {str(e)}")
    
    # Handle different webhook events
    event_type = payload.get("type")
    
    if event_type == "user.created":
        # A new user was created in Clerk
        user_data = payload.get("data", {})
        clerk_id = user_data.get("id")
        
        if not clerk_id:
            raise HTTPException(status_code=400, detail="Missing clerk_id in payload")
        
        # Check if user already exists
        existing_user = db.query(models.User).filter(models.User.clerk_id == clerk_id).first()
        if existing_user:
            return {"status": "User already exists"}
        
        # Create new user
        email = user_data.get("email_addresses", [{}])[0].get("email_address")
        first_name = user_data.get("first_name") 
        last_name = user_data.get("last_name")
        
        new_user = models.User(
            clerk_id=clerk_id,
            email=email,
            first_name=first_name,
            last_name=last_name,
            image_url=user_data.get("image_url")
        )
        
        db.add(new_user)
        db.commit()
        
        return {"status": "success", "message": "User created"}
        
    elif event_type == "user.updated":
        # User information was updated in Clerk
        user_data = payload.get("data", {})
        clerk_id = user_data.get("id")
        
        if not clerk_id:
            raise HTTPException(status_code=400, detail="Missing clerk_id in payload")
        
        # Find the user
        user = db.query(models.User).filter(models.User.clerk_id == clerk_id).first()
        if not user:
            raise HTTPException(status_code=404, detail="User not found")
        
        # Update user information
        email = user_data.get("email_addresses", [{}])[0].get("email_address")
        if email:
            user.email = email
        
        if user_data.get("first_name"):
            user.first_name = user_data.get("first_name")
        
        if user_data.get("last_name"):
            user.last_name = user_data.get("last_name")
        
        if user_data.get("image_url"):
            user.image_url = user_data.get("image_url")
        
        db.commit()
        return {"status": "success", "message": "User updated"}
        
    elif event_type == "user.deleted":
        # User was deleted in Clerk
        user_data = payload.get("data", {})
        clerk_id = user_data.get("id")
        
        if not clerk_id:
            raise HTTPException(status_code=400, detail="Missing clerk_id in payload")
        
        # Find the user
        user = db.query(models.User).filter(models.User.clerk_id == clerk_id).first()
        if not user:
            return {"status": "User not found"}
        
        # You can either delete the user or mark them as inactive
        # Here we'll delete the user
        db.delete(user)
        db.commit()
        
        return {"status": "success", "message": "User deleted"}
    
    # Return for other event types
    return {"status": "success", "message": f"Event {event_type} processed"}

@app.get("/auth-debug")
async def auth_debug(request: Request, current_user = Depends(get_current_user)):
    """Debug endpoint to check authentication"""
    print("==== Auth Debug Endpoint ====")
    
    # Print all headers for debugging
    print("All request headers:")
    for header_name, header_value in request.headers.items():
        print(f"  {header_name}: {header_value[:30] if len(header_value) > 30 else header_value}")
    
    # Check for auth header directly
    auth_header = request.headers.get("authorization") or request.headers.get("Authorization")
    
    if current_user:
        print(f"User authenticated via dependency: {current_user.id}, {current_user.clerk_id}")
        return {
            "authenticated": True,
            "user_id": current_user.id,
            "clerk_id": current_user.clerk_id,
            "email": current_user.email,
            "auth_header_present": auth_header is not None
        }
    elif auth_header:
        print(f"Auth header present but no user: {auth_header[:30]}...")
        return {
            "authenticated": False,
            "auth_header_present": True,
            "auth_header_prefix": auth_header[:30],
            "message": "Token received but couldn't authenticate user"
        }
    else:
        print("User not authenticated, no auth header")
        return {
            "authenticated": False,
            "auth_header_present": False
        }